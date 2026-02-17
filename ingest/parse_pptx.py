from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ingest.models import ExtractedUnit
from ingest.ocr import run_ocr, should_ocr_image


def _slide_native_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            chunks.append(shape.text.strip())
    return "\n".join(t for t in chunks if t)


def _slide_image_blobs(slide) -> list[bytes]:
    images: list[bytes] = []
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
            images.append(shape.image.blob)
    return images


def parse_pptx(path: Path, text_min_chars: int) -> tuple[list[ExtractedUnit], int]:
    prs = Presentation(path)
    units: list[ExtractedUnit] = []
    ocr_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        native_text = _slide_native_text(slide)
        if len(native_text) >= text_min_chars:
            units.append(
                ExtractedUnit(
                    source_type="ppt_slide",
                    page_or_slide=idx,
                    text=native_text,
                    text_origin="native_text",
                )
            )
            continue

        ocr_fragments: list[str] = []
        for blob in _slide_image_blobs(slide):
            if not should_ocr_image(blob):
                continue
            extracted = run_ocr(blob)
            if extracted:
                ocr_fragments.append(extracted)

        if ocr_fragments:
            units.append(
                ExtractedUnit(
                    source_type="ppt_slide",
                    page_or_slide=idx,
                    text="\n".join(ocr_fragments),
                    text_origin="ocr",
                )
            )
            ocr_count += 1
        elif native_text:
            units.append(
                ExtractedUnit(
                    source_type="ppt_slide",
                    page_or_slide=idx,
                    text=native_text,
                    text_origin="native_text",
                )
            )

    return units, ocr_count
